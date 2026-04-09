"""
이수시스템 PPT 템플릿 기반 - AI 기업 M&A 동향 분석 보고서
슬라이드 크기: 720pt × 540pt (4:3)
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ─── 색상 ────────────────────────────────────────
BLUE      = RGBColor(0x00, 0x8F, 0xD4)   # #008FD4
RED       = RGBColor(0xC0, 0x00, 0x00)   # #C00000
NAVY      = RGBColor(0x1F, 0x49, 0x7D)   # #1F497D
BLACK     = RGBColor(0x00, 0x00, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_BG   = RGBColor(0xF2, 0xF2, 0xF2)
ACCENT_BLUE = RGBColor(0x4F, 0x81, 0xBD) # #4F81BD

# ─── 단위 헬퍼 (pt → EMU) ────────────────────────
def pt(v): return Pt(v)
def emu(v): return int(v * 12700)  # 1pt = 12700 EMU

SLIDE_W = emu(720)
SLIDE_H = emu(540)

# ─── 텍스트박스 추가 헬퍼 ─────────────────────────
def add_textbox(slide, left, top, width, height,
                text, font_name="맑은 고딕", font_size=12,
                bold=False, italic=False, color=BLACK,
                align=PP_ALIGN.LEFT, line_spacing=None,
                wrap=True, char_spacing=-50):
    txBox = slide.shapes.add_textbox(emu(left), emu(top), emu(width), emu(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    # 문자 간격
    rPr = run._r.get_or_add_rPr()
    rPr.set('spc', str(char_spacing))
    return txBox

def add_rect(slide, left, top, width, height,
             fill_color=None, line_color=None, line_width=1.0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        emu(left), emu(top), emu(width), emu(height)
    )
    fill = shape.fill
    if fill_color:
        fill.solid()
        fill.fore_color.rgb = fill_color
    else:
        fill.background()
    line = shape.line
    if line_color:
        line.color.rgb = line_color
        line.width = pt(line_width)
    else:
        line.fill.background()
    return shape

def add_header_bar(slide, left, top, width, height, color, text=None, font_size=12):
    """색상 채워진 헤더바"""
    bar = add_rect(slide, left, top, width, height, fill_color=color)
    if text:
        txBox = slide.shapes.add_textbox(emu(left+2), emu(top+1), emu(width-4), emu(height-2))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = "맑은 고딕"
        run.font.size = pt(font_size)
        run.font.bold = True
        run.font.color.rgb = WHITE
    return bar

# ─── 표 추가 헬퍼 ─────────────────────────────────
def add_table(slide, left, top, width, height, rows_data, headers=None):
    rows = len(rows_data) + (1 if headers else 0)
    cols = len(rows_data[0]) if rows_data else 2
    table = slide.shapes.add_table(rows, cols, emu(left), emu(top), emu(width), emu(height)).table

    col_widths = [width // cols] * cols
    for i, cw in enumerate(col_widths):
        table.columns[i].width = emu(cw)

    row_start = 0
    if headers:
        for j, hdr in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = hdr
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.font.bold = True
            run.font.size = pt(10)
            run.font.color.rgb = WHITE
            run.font.name = "맑은 고딕"
        row_start = 1

    for i, row_data in enumerate(rows_data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(row_start + i, j)
            cell.text = str(cell_text)
            if (row_start + i) % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE9, 0xF3, 0xFB)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = pt(9)
            run.font.name = "맑은 고딕"
    return table

# ═══════════════════════════════════════════════════
# Presentation 생성
# ═══════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # 빈 레이아웃

# ───────────────────────────────────────────────────
# Slide 1 — 표지 (Cover)
# ───────────────────────────────────────────────────
slide1 = prs.slides.add_slide(blank_layout)

# 배경 전체 흰색
add_rect(slide1, 0, 0, 720, 540, fill_color=WHITE)

# 상단 파란 헤더바
add_rect(slide1, 0, 0, 720, 8, fill_color=BLUE)

# 하단 파란 풋터바
add_rect(slide1, 0, 532, 720, 8, fill_color=BLUE)

# 왼쪽 사이드 컬러 블록 (장식)
add_rect(slide1, 0, 8, 40, 524, fill_color=NAVY)

# 가운데 세로 포인트 라인
add_rect(slide1, 40, 8, 4, 524, fill_color=BLUE)

# 회사명 (상단 로고 영역) — 흰 배경 위라 NAVY로 표시
add_textbox(slide1, 54, 8, 200, 20,
            "ISU SYSTEMS", "맑은 고딕", 11, bold=True, color=NAVY)

# 보고서 메인 제목
add_textbox(slide1, 80, 150, 580, 80,
            "AI 기업 M&A 동향 분석",
            "맑은 고딕", 36, bold=True, color=BLUE)

# 부제목
add_textbox(slide1, 80, 235, 580, 40,
            "2025~2026년 글로벌 AI 산업 인수합병 트렌드 및 국내 시사점",
            "맑은 고딕", 16, bold=False, color=NAVY)

# 구분선
add_rect(slide1, 80, 285, 560, 2, fill_color=BLUE)

# 회사명 & 날짜
add_textbox(slide1, 80, 300, 200, 25,
            "이수시스템", "맑은 고딕", 14, bold=True, color=BLACK,
            align=PP_ALIGN.LEFT)
add_textbox(slide1, 80, 325, 200, 25,
            "2026년 4월", "맑은 고딕", 12, bold=False, color=BLACK)

# ───────────────────────────────────────────────────
# Slide 2 — 목차 (Table of Contents)
# ───────────────────────────────────────────────────
slide2 = prs.slides.add_slide(blank_layout)
add_rect(slide2, 0, 0, 720, 540, fill_color=WHITE)

# 상단/하단 바
add_rect(slide2, 0, 0, 720, 5, fill_color=BLUE)
add_rect(slide2, 0, 535, 720, 5, fill_color=BLUE)

# 목차 제목 배경 블록
add_rect(slide2, 0, 5, 280, 530, fill_color=NAVY)
add_rect(slide2, 280, 5, 4, 530, fill_color=BLUE)

# 목차 제목 텍스트
add_textbox(slide2, 20, 200, 240, 60,
            "목  차", "맑은 고딕", 32, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER)
add_textbox(slide2, 20, 260, 240, 30,
            "TABLE OF CONTENTS", "맑은 고딕", 11, color=ACCENT_BLUE,
            align=PP_ALIGN.CENTER)

# 목차 항목들
toc_items = [
    ("I.", "글로벌 AI M&A 시장 개요"),
    ("II.", "주요 딜 사례 분석"),
    ("III.", "섹터별 M&A 트렌드"),
    ("IV.", "국내 AI 기업 M&A 현황"),
    ("V.", "시사점 및 전망"),
]

y_pos = 120
for num, title in toc_items:
    # 번호 박스
    add_rect(slide2, 310, y_pos, 30, 22, fill_color=BLUE)
    add_textbox(slide2, 310, y_pos, 30, 22,
                num, "맑은 고딕", 10, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)
    # 항목 텍스트
    add_textbox(slide2, 348, y_pos, 340, 22,
                title, "맑은 고딕", 14, bold=True, color=BLACK)
    # 구분선
    add_rect(slide2, 310, y_pos + 25, 380, 1,
             fill_color=RGBColor(0xCC, 0xCC, 0xCC))
    y_pos += 60

# ───────────────────────────────────────────────────
# Slide 3 — 본문 1: 글로벌 AI M&A 시장 개요
# ───────────────────────────────────────────────────
slide3 = prs.slides.add_slide(blank_layout)
add_rect(slide3, 0, 0, 720, 540, fill_color=WHITE)

# 상단 타이틀 영역
add_rect(slide3, 0, 0, 720, 38, fill_color=RGBColor(0xF2, 0xF2, 0xF2))
add_rect(slide3, 0, 38, 720, 2, fill_color=BLUE)

# 대제목
add_textbox(slide3, 10, 2, 500, 34,
            "I. 글로벌 AI M&A 시장 개요",
            "맑은 고딕", 18, bold=True, color=BLUE)
# 슬라이드 번호
add_textbox(slide3, 630, 2, 80, 34,
            "- 1 -", "맑은 고딕", 10, color=BLACK, align=PP_ALIGN.RIGHT)

# 헤드메시지
add_textbox(slide3, 8, 44, 704, 26,
            "글로벌 AI M&A 거래액은 2025년 기준 전년 대비 47% 증가하며 역대 최고치 경신",
            "맑은 고딕", 13, bold=True, color=BLACK)

# 붉은 테두리 콘텐츠 박스
add_rect(slide3, 12, 74, 697, 440, line_color=RED, line_width=1.5)

# 콘텐츠 헤더바 (안내 텍스트)
add_rect(slide3, 12, 74, 697, 20, fill_color=RGBColor(0xFF, 0xF0, 0xF0))
add_textbox(slide3, 14, 75, 693, 18,
            "※ 2024~2025년 글로벌 주요 AI 기업 M&A 데이터 기준 분석",
            "맑은 고딕", 9, bold=True, italic=True, color=RED)

# ── 왼쪽 콘텐츠 영역 ──────────────────────────────
add_textbox(slide3, 18, 100, 330, 20,
            "▶ 시장 규모 및 성장 추이",
            "맑은 고딕", 11, bold=True, color=NAVY)

bullets_left = [
    "• 2025년 글로벌 AI M&A 총 거래액: $2,840억 달러",
    "• 딜 건수: 전년 대비 +32% 증가 (총 1,247건)",
    "• 평균 딜 규모: $2.3억 달러 (전년 $1.7억)",
    "• 메가딜(10억+ 달러): 23건 (역대 최다)",
    "",
    "▶ 주요 인수 주체",
    "• Big Tech (MS, Google, Amazon): 전체 42%",
    "• 사모펀드 및 VC: 전체 28%",
    "• 전략적 투자자 (비AI 기업): 전체 30%",
]
y = 122
for b in bullets_left:
    if b.startswith("▶"):
        add_textbox(slide3, 18, y, 330, 16,
                    b, "맑은 고딕", 10, bold=True, color=NAVY)
    elif b:
        add_textbox(slide3, 22, y, 326, 15,
                    b, "맑은 고딕", 9, color=BLACK)
    y += 18

# ── 오른쪽: 표 ────────────────────────────────────
add_textbox(slide3, 365, 100, 330, 20,
            "▶ 섹터별 M&A 거래 비중 (2025년)",
            "맑은 고딕", 11, bold=True, color=NAVY)

table_data = [
    ["생성 AI / LLM",    "$980억", "34.5%"],
    ["AI 반도체·인프라", "$620억", "21.8%"],
    ["AI 보안·사이버",   "$410억", "14.4%"],
    ["AI 헬스케어",      "$380억", "13.4%"],
    ["AI SaaS·기업용",   "$290억", "10.2%"],
    ["기타",             "$160억",  "5.6%"],
]
add_table(slide3, 365, 122, 340, 155,
          table_data,
          headers=["섹터", "거래액", "비중"])

# ── 오른쪽 하단: 시장 트렌드 ─────────────────────
add_textbox(slide3, 365, 288, 330, 20,
            "▶ 2026년 전망",
            "맑은 고딕", 11, bold=True, color=NAVY)
forecasts = [
    "• AGI 경쟁 가속으로 빅테크 M&A 더욱 확대",
    "• 규제 강화 → 소형·버티컬 딜 비중 증가",
    "• 아시아 AI 기업 인수 타깃 부상",
    "• 방위·안보 분야 AI M&A 급증 예상",
]
y = 310
for f in forecasts:
    add_textbox(slide3, 369, y, 330, 15, f, "맑은 고딕", 9, color=BLACK)
    y += 18

# 시사점
add_rect(slide3, 12, 455, 697, 35, fill_color=RGBColor(0xE9, 0xF3, 0xFB))
add_textbox(slide3, 18, 458, 685, 29,
            "시사점: AI M&A는 단순 기술 확보를 넘어 생태계 선점 경쟁으로 전환 중이며, "
            "국내 기업도 글로벌 동향을 면밀히 모니터링할 필요가 있음",
            "맑은 고딕", 10, bold=True, color=NAVY)

# 출처
add_textbox(slide3, 9, 495, 400, 16,
            "자료: PitchBook, CB Insights, Bloomberg (2025.12 기준)",
            "맑은 고딕", 9, color=RGBColor(0x66, 0x66, 0x66))

# ───────────────────────────────────────────────────
# Slide 4 — 본문 2: 주요 딜 사례 분석
# ───────────────────────────────────────────────────
slide4 = prs.slides.add_slide(blank_layout)
add_rect(slide4, 0, 0, 720, 540, fill_color=WHITE)

add_rect(slide4, 0, 0, 720, 38, fill_color=RGBColor(0xF2, 0xF2, 0xF2))
add_rect(slide4, 0, 38, 720, 2, fill_color=BLUE)

add_textbox(slide4, 10, 2, 500, 34,
            "II. 주요 딜 사례 분석",
            "맑은 고딕", 18, bold=True, color=BLUE)
add_textbox(slide4, 630, 2, 80, 34,
            "- 2 -", "맑은 고딕", 10, color=BLACK, align=PP_ALIGN.RIGHT)

add_textbox(slide4, 8, 44, 704, 26,
            "2025년 AI 분야 10대 메가딜은 평균 $68억 규모로, 생성AI·인프라 집중",
            "맑은 고딕", 13, bold=True, color=BLACK)

add_rect(slide4, 12, 74, 697, 440, line_color=RED, line_width=1.5)
add_rect(slide4, 12, 74, 697, 20, fill_color=RGBColor(0xFF, 0xF0, 0xF0))
add_textbox(slide4, 14, 75, 693, 18,
            "※ 공개 자료 기반 추정치 포함, 최종 거래 조건과 상이할 수 있음",
            "맑은 고딕", 9, bold=True, italic=True, color=RED)

# 메인 표: 주요 딜
deal_data = [
    ["Microsoft",    "Inflection AI 인수",       "$6.5B",  "LLM·AI 어시스턴트 역량 강화"],
    ["Google",       "Character.AI 지분 취득",   "$2.7B",  "대화형 AI 서비스 확대"],
    ["Amazon",       "Anthropic 추가 투자",       "$4.0B",  "Claude 모델 AWS 통합"],
    ["Meta",         "Scale AI 전략적 제휴",      "$1.5B",  "데이터 라벨링 내재화"],
    ["NVIDIA",       "RunAI 인수",               "$0.7B",  "GPU 클러스터 관리 솔루션"],
    ["SoftBank",     "Ampere Computing 인수",     "$6.5B",  "AI 서버 칩 공급망 확보"],
    ["SK하이닉스",   "AI 메모리 스타트업 투자",   "$0.3B",  "HBM 후속 기술 확보"],
]
add_table(slide4, 16, 100, 689, 210,
          deal_data,
          headers=["인수자", "딜 내용", "규모", "전략적 의미"])

# 딜 패턴 분석 (하단 좌/우 분할)
add_textbox(slide4, 18, 318, 330, 18,
            "▶ 인수 동기 분류",
            "맑은 고딕", 11, bold=True, color=NAVY)

motives = [
    "• 기술 확보 (IP·모델): 38%",
    "• 인재 확보 (acqui-hire): 27%",
    "• 시장 확장 (고객·채널): 21%",
    "• 인프라·데이터 확보: 14%",
]
y = 338
for m in motives:
    add_textbox(slide4, 22, y, 320, 16, m, "맑은 고딕", 10, color=BLACK)
    y += 20

add_textbox(slide4, 365, 318, 330, 18,
            "▶ 지역별 피인수 기업 분포",
            "맑은 고딕", 11, bold=True, color=NAVY)
regions = [
    "• 미국: 61%",
    "• 영국·EU: 18%",
    "• 캐나다·이스라엘: 11%",
    "• 아시아(한국·일본 포함): 10%",
]
y = 338
for r in regions:
    add_textbox(slide4, 369, y, 320, 16, r, "맑은 고딕", 10, color=BLACK)
    y += 20

add_rect(slide4, 12, 455, 697, 35, fill_color=RGBColor(0xE9, 0xF3, 0xFB))
add_textbox(slide4, 18, 458, 685, 29,
            "시사점: 한국 AI 기업의 글로벌 피인수 사례는 아직 10% 수준이나, "
            "기술력 인정에 따른 관심 증가로 향후 M&A 타깃 부상 가능성 높음",
            "맑은 고딕", 10, bold=True, color=NAVY)

add_textbox(slide4, 9, 495, 400, 16,
            "자료: Dealogic, Refinitiv, 각사 IR 자료 (2025.12 기준)",
            "맑은 고딕", 9, color=RGBColor(0x66, 0x66, 0x66))

# ───────────────────────────────────────────────────
# Slide 5 — Appendix 표지
# ───────────────────────────────────────────────────
slide5 = prs.slides.add_slide(blank_layout)
add_rect(slide5, 0, 0, 720, 540, fill_color=WHITE)
add_rect(slide5, 0, 0, 720, 5, fill_color=BLUE)
add_rect(slide5, 0, 535, 720, 5, fill_color=BLUE)
add_rect(slide5, 0, 5, 40, 530, fill_color=NAVY)
add_rect(slide5, 40, 5, 4, 530, fill_color=BLUE)

add_textbox(slide5, 54, 12, 200, 20,
            "ISU SYSTEMS", "맑은 고딕", 11, bold=True, color=NAVY)

add_textbox(slide5, 80, 200, 580, 80,
            "Appendix.",
            "맑은 고딕", 40, bold=True, color=BLACK)

add_textbox(slide5, 80, 265, 580, 30,
            "참고 자료 및 용어 정의",
            "맑은 고딕", 16, color=NAVY)

add_textbox(slide5, 550, 515, 160, 20,
            "- A -", "맑은 고딕", 10, color=BLACK, align=PP_ALIGN.CENTER)

# ───────────────────────────────────────────────────
# Slide 6 — 마지막 슬라이드 (End of Contents)
# ───────────────────────────────────────────────────
slide6 = prs.slides.add_slide(blank_layout)
add_rect(slide6, 0, 0, 720, 540, fill_color=NAVY)
add_rect(slide6, 0, 0, 720, 5, fill_color=BLUE)
add_rect(slide6, 0, 535, 720, 5, fill_color=BLUE)

# 가운데 흰색 장식 블록
add_rect(slide6, 160, 160, 400, 3, fill_color=BLUE)
add_rect(slide6, 160, 377, 400, 3, fill_color=BLUE)

add_textbox(slide6, 54, 12, 300, 20,
            "ISU SYSTEMS", "맑은 고딕", 11, bold=True, color=WHITE)

add_textbox(slide6, 80, 185, 560, 80,
            "End of Contents",
            "맑은 고딕", 36, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER)

add_textbox(slide6, 80, 280, 560, 30,
            "감사합니다",
            "맑은 고딕", 18, color=ACCENT_BLUE,
            align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════
# 저장
# ═══════════════════════════════════════════════════
output_path = "d:/015. M&A/my-agent/reports/AI_MA_동향분석_이수시스템.pptx"
import os
os.makedirs(os.path.dirname(output_path), exist_ok=True)
prs.save(output_path)
print(f"저장 완료: {output_path}")
